"""
Rebuild defense PPT with architecture diagram slide at position 7.
Combines: original defense slides + arch diagram slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from copy import deepcopy

base = Path(__file__).resolve().parent.parent
defense_path = base / "docs" / "superpowers" / "specs" / "defense_ppt.pptx"
arch_path = base / "docs" / "superpowers" / "specs" / "arch_diagram.pptx"

# Load both
prs_def = Presentation(str(defense_path))
prs_arch = Presentation(str(arch_path))

arch_slide = prs_arch.slides[0]
total = len(prs_def.slides)
print(f"Defense PPT has {total} slides")

# Strategy: delete old slide 6 (index 6, 0-based), the old arch slide
# Then insert the new arch slide at that position

# First, remove all shapes from slide 6 and copy arch shapes onto it
target_slide = prs_def.slides[6]  # index 6 = slide 7 (old architecture placeholder)

# Remove existing shapes from target slide (except background)
shapes_to_remove = list(target_slide.shapes)
for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Copy background from arch slide
# Get arch slide background
arch_bg = arch_slide.background
target_bg = target_slide.background
# Copy fill
if arch_bg.fill.type is not None:
    target_bg.fill.solid()
    try:
        target_bg.fill.fore_color.rgb = arch_bg.fill.fore_color.rgb
    except:
        pass

# Copy all shapes from arch slide to target slide
for shape in arch_slide.shapes:
    el = deepcopy(shape._element)
    target_slide.shapes._spTree.append(el)

prs_def.save(str(defense_path))
print(f"Done! Architecture diagram inserted at slide 7. Total: {len(prs_def.slides)} slides")
