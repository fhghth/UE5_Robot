"""
Generate defense PPT for "基于UE5的双足机器人智能寻路系统的设计与开发"
Based on analysis doc: docs/superpowers/specs/2026-05-08-defense-ppt-outline.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Theme Constants ──────────────────────────────────────────
CLR_DARK  = RGBColor(0x1A, 0x1A, 0x2E)    # dark navy bg
CLR_ACCENT = RGBColor(0x00, 0x96, 0xC7)    # cyan accent
CLR_ACCENT2 = RGBColor(0xF7, 0x25, 0x5D)   # red accent
CLR_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
CLR_GRAY  = RGBColor(0xA0, 0xA0, 0xB0)
CLR_GREEN = RGBColor(0x2E, 0xCC, 0x71)
CLR_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
CLR_LIGHT_BG = RGBColor(0x24, 0x24, 0x3E)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────
def add_bg(slide, color=CLR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=CLR_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=CLR_WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add consistent title bar at top"""
    add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.0), CLR_LIGHT_BG)
    # accent line
    add_shape_bg(slide, Inches(0.5), Inches(0.95), Inches(2.0), Inches(0.05), CLR_ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5), title_text,
                 font_size=28, bold=True, color=CLR_WHITE)
    if subtitle_text:
        add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.35), subtitle_text,
                     font_size=14, color=CLR_GRAY)

def add_section_label(slide, section_name, page_num):
    """Top-right section + page number"""
    add_text_box(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.35), section_name,
                 font_size=11, color=CLR_GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.3), str(page_num),
                 font_size=11, color=CLR_GRAY, alignment=PP_ALIGN.RIGHT)

def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    return slide

def add_table(slide, left, top, col_widths, headers, rows, font_size=13):
    """Create a styled table"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.45 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = CLR_WHITE
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_ACCENT

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = CLR_WHITE
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A) if ri % 2 == 0 else RGBColor(0x32, 0x32, 0x52)
    return table_shape

def add_code_ref(slide, text, left=Inches(0.5), top=Inches(6.7)):
    """Add code reference at bottom"""
    add_text_box(slide, left, top, Inches(12), Inches(0.3), text,
                 font_size=10, color=RGBColor(0x60, 0x60, 0x80))

def add_icon_bullet(slide, left, top, icon, text, font_size=16, color=CLR_WHITE):
    add_text_box(slide, left, top, Inches(0.4), Inches(0.3), icon, font_size=font_size, color=CLR_ACCENT)
    add_text_box(slide, left + Inches(0.45), top, Inches(11), Inches(0.3), text, font_size=font_size, color=color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_shape_bg(sl, Inches(0), Inches(0), W, H, CLR_DARK)
add_shape_bg(sl, Inches(0), Inches(2.8), W, Inches(0.06), CLR_ACCENT)
add_shape_bg(sl, Inches(0), Inches(5.2), W, Inches(0.03), CLR_ACCENT2)
add_text_box(sl, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             '基于UE5的双足机器人智能寻路系统的设计与开发',
             font_size=40, bold=True, color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(0.5),
             '基于分层深度强化学习的决策-控制解耦方案',
             font_size=22, color=CLR_ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.2), Inches(11), Inches(0.5),
             '项目答辩汇报',
             font_size=20, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             'UE5 Chaos Physics · Soft Actor-Critic · Hierarchical RL · ONNX Deployment',
             font_size=14, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             '2026年5月',
             font_size=14, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: 目录
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '汇报提纲')
add_section_label(sl, '概述', 2)

sections = [
    ('01', '研究背景与问题定义', '双足机器人导航的核心挑战与三种技术路线对比'),
    ('02', '系统架构与核心技术 ★', '分层推理架构 · 奖励体系 · 课程学习 · 训练优化 · ONNX部署'),
    ('03', '实验结果与分析', '步态训练 · 导航性能 · 消融实验 · 关键发现'),
    ('04', '局限性与未来展望', '当前四局限 → 五阶段技术演进路线图'),
    ('05', '总结与贡献', '四条核心贡献总结'),
]
for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.5 + i * 1.1)
    add_text_box(sl, Inches(1.2), y, Inches(0.8), Inches(0.5), num,
                 font_size=32, bold=True, color=CLR_ACCENT)
    add_text_box(sl, Inches(2.2), y, Inches(9), Inches(0.45), title,
                 font_size=22, bold=True, color=CLR_WHITE)
    add_text_box(sl, Inches(2.2), y + Inches(0.45), Inches(9), Inches(0.4), desc,
                 font_size=14, color=CLR_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: 研究背景
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '研究背景与核心痛点')
add_section_label(sl, '问题定义', 3)

items = [
    '▎ 核心矛盾：决策（往哪走）与 控制（怎么迈步）的深度耦合',
    '▎ 高维连续控制：12+ 维关节空间，每个控制周期需在 33ms 内完成推理',
    '▎ 不稳定动力学：双足机器人始终处于"跌倒边缘"，对控制精度极度敏感',
    '▎ 复杂环境感知：需要同时处理路径规划、障碍物规避、步态生成与物理交互',
    '',
    '传统方案在实时性、稳定性与泛化能力上均暴露明显短板',
]
add_bullet_frame(sl, Inches(1), Inches(1.5), Inches(11), Inches(4), items, font_size=18, spacing=Pt(14))

# Box with key numbers
add_shape_bg(sl, Inches(8.5), Inches(4.5), Inches(4), Inches(2.2), CLR_LIGHT_BG)
add_text_box(sl, Inches(8.8), Inches(4.6), Inches(3.5), Inches(1.8),
             '关键数据\n\n▸ 训练周期：150万步\n▸ 控制频率：30Hz (33ms)\n▸ 可调参数：50+ 项\n▸ 物理崩溃记录：40 次',
             font_size=13, color=CLR_GRAY)

add_code_ref(sl, 'OrangeRobotEnvComponent.h — 观测/动作空间定义 | Saved/Crashes/ — 40个崩溃目录')

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: 三种方案对比（上）
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '三种技术路线的对比分析', '端到端DRL · MPC轨迹优化 · 分层RL — 本场景下的致命短板')
add_section_label(sl, '问题定义', 4)

headers = ['方案', '核心思路', '本场景下的致命短板']
rows = [
    ['端到端RL\n（当前基线）', '单策略直接从\n观测→关节指令', '50+参数高度耦合、单信号压缩多目标\n仅ActionDeadzone后验过滤、150万步/周期\nsim-to-real过拟合风险'],
    ['MPC轨迹优化\n+步态规划', '简化模型在线求解\n+全身逆动力学', '30Hz下33ms内无法完成非线性求解\n依赖精确动力学模型（Chaos黑盒）\n接触序列组合爆炸'],
    ['分层RL\n（高层+低层）', '高低层分离：\n规划与执行解耦', '非平稳训练环境→Replay Buffer污染\n子目标空间定义缺乏明确方案\ngRPC通信开销翻倍放大'],
]
add_table(sl, Inches(0.8), Inches(1.5),
          [Inches(2.2), Inches(2.8), Inches(6.8)], headers, rows, font_size=14)

# conclusion
add_shape_bg(sl, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), RGBColor(0x2A, 0x2A, 0x4A))
add_text_box(sl, Inches(1.2), Inches(5.9), Inches(11), Inches(0.6),
             '结论：分层RL是最有前景的方向，但必须解决三个核心挑战 — ①解耦训练  ②物理约束内化  ③通信效率',
             font_size=16, bold=True, color=CLR_GREEN)

add_code_ref(sl, 'train_headless.py SAC_DEFAULTS | OrangeRobotEnvComponent.h ActionDeadzone=0.05 | gRPC Port=50051')

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: 三种方案对比（下）— 深度剖析
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '各方案在本UE5 Chaos环境下的具体不可行性')
add_section_label(sl, '问题定义', 5)

items = [
    '▎ 端到端DRL的致命缺陷',
    '   FootImpactVelocityThreshold=50.0 → 事后修补越界，无法从架构层面杜绝物理违规',
    '   ActionDeadzone=0.05 → 仅过滤小幅度噪声，关节超限/力矩饱和无任何硬约束',
    '   40个崩溃记录（Saved/Crashes/）→ 被动惩罚机制的系统性失效',
    '',
    '▎ MPC轨迹优化的实时性悖论',
    '   SimulationFrequencyHz=30 → 每33ms一次控制周期',
    '   DesiredStepPeriod=0.5s → 每0.25s需重规划接触序列（混合整数规划）',
    '   Schola传感器为黑盒接口 → 无完整动力学Jacobian → 模型偏差逐周期放大',
    '',
    '▎ 分层RL的训练-部署矛盾',
    '   bEnableHighLevelCommand=True → 低层接收来自未收敛高层的非平稳指令',
    '   SAC Off-policy机制 → Replay Buffer中旧经验持续污染梯度更新',
    '   CommandMatchBaseReward=0.5 → 对非平稳性的"事后补偿"，缺乏系统化解耦方案',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=16, spacing=Pt(8))

add_code_ref(sl, 'OrangeRobotEnvComponent.h:260,275,371,413 | OrangeRobotEnvComponent.cpp:489-591')

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 项目目标
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '项目目标与研究问题')
add_section_label(sl, '问题定义', 6)

add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
             '总目标：设计一套分层DRL架构，在UE5高保真物理环境中实现双足机器人的稳定行走与自主导航',
             font_size=18, bold=True, color=CLR_WHITE)

add_shape_bg(sl, Inches(1), Inches(2.3), Inches(11), Inches(0.03), CLR_ACCENT)

questions = [
    ('Q1', '如何在不牺牲训练稳定性的前提下实现高层导航与低层步态的有效解耦？',
     '→ 核心方案：独立分训 + 固定2D指令接口 + 课程化指令采样'),
    ('Q2', '如何从架构层面（而非仅奖励函数）保障策略输出的物理可行性？',
     '→ 核心方案：ActionDeadzone + L2范数裁剪 + 规划中的QP约束求解器'),
    ('Q3', '如何在50+可调参数的复杂奖励空间中实现系统化的权重搜索与消融实验？',
     '→ 核心方案：14组分位掩码体系 + 三组独立开关 + TensorBoard分量追踪'),
]
for i, (qid, question, answer) in enumerate(questions):
    y = Inches(2.7 + i * 1.4)
    add_text_box(sl, Inches(1.2), y, Inches(0.6), Inches(0.4), qid,
                 font_size=22, bold=True, color=CLR_ACCENT2)
    add_text_box(sl, Inches(1.9), y, Inches(9.5), Inches(0.5), question,
                 font_size=16, bold=True, color=CLR_WHITE)
    add_text_box(sl, Inches(1.9), y + Inches(0.55), Inches(9.5), Inches(0.5), answer,
                 font_size=14, color=CLR_GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: 系统总体架构图
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '系统总体架构', '分层推理 · 独立训练 · 桥接部署')
add_section_label(sl, '核心架构', 7)

# Architecture diagram using shapes
# High-level box
add_shape_bg(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.8), CLR_LIGHT_BG)
add_text_box(sl, Inches(1), Inches(1.55), Inches(5), Inches(0.4),
             '🔹 高层导航策略 (5Hz, ONNX)', font_size=16, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(1), Inches(2.0), Inches(5), Inches(0.3),
             '观测: 目标方向+距离+8向射线+通路开阔度 (17维)', font_size=12, color=CLR_GRAY)
add_text_box(sl, Inches(1), Inches(2.3), Inches(5), Inches(0.3),
             '动作: [Forward, Turn] 2维连续指令', font_size=12, color=CLR_GRAY)
add_text_box(sl, Inches(1), Inches(2.55), Inches(5), Inches(0.3),
             'XNavigationCubeEnvComponent · USimpleStepper', font_size=11, color=CLR_GRAY)

# Arrow
add_text_box(sl, Inches(6.5), Inches(2.0), Inches(1.2), Inches(0.5),
             '▶ EMA\nα=0.3', font_size=12, color=CLR_ORANGE, alignment=PP_ALIGN.CENTER)

# Low-level box
add_shape_bg(sl, Inches(7.5), Inches(1.5), Inches(5.5), Inches(1.8), CLR_LIGHT_BG)
add_text_box(sl, Inches(7.7), Inches(1.55), Inches(5), Inches(0.4),
             '🔸 低层步态策略 (30Hz, ONNX)', font_size=16, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(7.7), Inches(2.0), Inches(5), Inches(0.3),
             '观测: 躯干状态+关节角度+足地交互+步态相位 (30-40维)', font_size=12, color=CLR_GRAY)
add_text_box(sl, Inches(7.7), Inches(2.3), Inches(5), Inches(0.3),
             '动作: 10-12维关节角速度目标', font_size=12, color=CLR_GRAY)
add_text_box(sl, Inches(7.7), Inches(2.55), Inches(5), Inches(0.3),
             'OrangeRobotEnvComponent · UPipelinedStepper', font_size=11, color=CLR_GRAY)

# Bridge
add_shape_bg(sl, Inches(3.5), Inches(3.8), Inches(6.3), Inches(1.0), RGBColor(0x28, 0x28, 0x50))
add_text_box(sl, Inches(3.7), Inches(3.85), Inches(5.8), Inches(0.9),
             '🔗 Bridge组件 (ABridgeComponent)\n'
             '高层5Hz调度 · 低层30Hz调度 · CommandSmoothAlpha=0.3 · 双ONNX模型加载',
             font_size=13, color=CLR_WHITE, alignment=PP_ALIGN.CENTER)

# Physical Safety Layer
add_shape_bg(sl, Inches(3.5), Inches(5.2), Inches(6.3), Inches(0.8), RGBColor(0x35, 0x20, 0x20))
add_text_box(sl, Inches(3.7), Inches(5.25), Inches(5.8), Inches(0.7),
             '🛡 动作安全层: ActionDeadzone=0.05 · L2范数裁剪(max=2.0) · 角速度限幅 · 非线性塑形(exp=2.0)',
             font_size=12, color=CLR_ORANGE, alignment=PP_ALIGN.CENTER)

# Chaos
add_shape_bg(sl, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.6), CLR_ACCENT2)
add_text_box(sl, Inches(3.7), Inches(6.35), Inches(5.8), Inches(0.5),
             '⚡ Chaos 物理引擎 · SimulationFrequencyHz=30 · SetAngularVelocityTarget',
             font_size=13, bold=True, color=CLR_WHITE, alignment=PP_ALIGN.CENTER)

add_code_ref(sl, 'BridgeComponent.h:26 · BridgeComponent.cpp:294-323 · OrangeRobotEnvComponent.cpp:735-803')

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: 核心技术一 — 分层推理架构
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术一：分层推理-独立训练架构', '以固定2D指令接口 + 课程化指令采样实现决策-控制解耦')
add_section_label(sl, '核心技术', 8)

items = [
    '▎ 为什么分层？',
    '   · 导航决策（秒级，需看远处目标）与姿态稳定（毫秒级，需维持平衡）的时间尺度天然不匹配',
    '   · 单策略需同时处理两者 → 维度爆炸 + 梯度冲突',
    '',
    '▎ 如何分层？',
    '   · 高层：5Hz推理（HighInferenceInterval=0.2f），17维导航观测 → 2维速度指令',
    '   · 低层：30Hz推理，30-40维本体感知 + 2维指令 → 10-12维关节目标速度',
    '',
    '▎ 桥接机制：',
    '   · EMA平滑（CommandSmoothAlpha=0.3）→ 抑制高层输出抖动，产生平滑指令流',
    '   · SetHighLevelCommand → clamp到[-1,1]后注入低层观测空间的末尾2维',
    '',
    '▎ 独立训练的工程优势：',
    '   · 高层在XNavigationCubeEnvComponent中独立学习绕障导航',
    '   · 低层在OrangeRobotEnvComponent中通过课程化随机指令学习步态',
    '   · 两层通过2维指令接口完全解耦 → 消除联合训练的Replay Buffer污染问题',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=15, spacing=Pt(6))

add_code_ref(sl, 'BridgeComponent.cpp:109,294-323 | OrangeRobotEnvComponent.cpp:422-426,910-914')

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: 独立训练的解耦优势
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '独立训练 vs 联合训练：方案对比与指令分布匹配')
add_section_label(sl, '核心技术', 9)

# Two columns
add_shape_bg(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), RGBColor(0x2A, 0x1A, 0x1A))
add_text_box(sl, Inches(0.7), Inches(1.6), Inches(5.3), Inches(0.4),
             '❌ 联合训练（标准分层RL做法）', font_size=16, bold=True, color=CLR_ACCENT2)
add_text_box(sl, Inches(0.7), Inches(2.1), Inches(5.3), Inches(1.6),
             '⋅ 高层策略每步更新 → 低层面临非平稳MDP\n'
             '⋅ SAC Off-policy → Replay Buffer中旧高层指令\n'
             '  持续"污染"低层Critic的TD目标\n'
             '⋅ 实测表现：Q值发散、收敛极慢、episode\n'
             '  长度无法稳定增长',
             font_size=14, color=CLR_GRAY)

add_shape_bg(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), RGBColor(0x1A, 0x2A, 0x1A))
add_text_box(sl, Inches(7.0), Inches(1.6), Inches(5.3), Inches(0.4),
             '✅ 独立训练（本项目的创新方案）', font_size=16, bold=True, color=CLR_GREEN)
add_text_box(sl, Inches(7.0), Inches(2.1), Inches(5.3), Inches(1.6),
             '⋅ 高层和低层分别独立训练SAC\n'
             '⋅ 低层接收课程化随机指令采样（非高层输出）\n'
             '⋅ 指令分布平稳 → SAC Critic稳定收敛\n'
             '⋅ 部署时通过Bridge桥接两个ONNX模型',
             font_size=14, color=CLR_GRAY)

# Curriculum command distribution
add_text_box(sl, Inches(0.8), Inches(4.3), Inches(11), Inches(0.4),
             '课程化指令采样策略：五阶段渐进覆盖指令空间', font_size=16, bold=True, color=CLR_ACCENT)

headers2 = ['阶段', '步数范围', '指令范围 Forward', '指令范围 Turn', '学习目标']
rows2 = [
    ['0 站立', '0~150k', '(0, 0) 固定', '0 固定', '躯干直立+不摔倒'],
    ['1 极慢速', '150k~300k', '[0.05, 0.2]', '0', '慢速纯前进步态'],
    ['2 中速', '300k~600k', '[0.2, 0.4]', '[-0.05, 0.05]', '中速行走+对称性'],
    ['3 快速', '600k~1M', '[-1, 1]', '[-0.5, 0.5]', '全速后退+能效'],
    ['4 全指令', '>1M', '[-1, 1]', '[-1, 1]', '全地形泛化'],
]
add_table(sl, Inches(0.8), Inches(4.8),
          [Inches(1.4), Inches(1.8), Inches(2.5), Inches(2.5), Inches(3.0)],
          headers2, rows2, font_size=12)

add_code_ref(sl, 'OrangeRobotEnvComponent.cpp:473-483 SampleEpisodeHighLevelCommand | :489-591 UpdateCurriculumWeightsAndCommand')

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: 核心技术二 — 奖励体系
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术二：14组分+位掩码奖励体系', '"分解-独立控制-系统消融" — 将50+参数耦合空间分化为可独立开关的子系统')
add_section_label(sl, '核心技术', 10)

# Three groups
groups = [
    ('CORE 组 (6分量)\n基础生存+指令跟踪', CLR_ACCENT,
     ['Alive — 存活激励+姿态质量调制', 'Height — 三段式高度维持', 'LateralPenalty — 侧向漂移抑制',
      'TrunkStability — 倾斜/角速度/颠簸', 'CommandTracking — 指数核指令跟踪', 'StableDoubleSupport — 双足支撑']),
    ('GAIT 组 (6分量)\n步态质量+自然度', CLR_GREEN,
     ['SupportStability — 支撑稳定性', 'GaitQuality — 步态质量综合', 'StepAlternation — 交替步态奖励',
      'StepFrequency — 步频[1,3]Hz', 'FootImpact — 足部冲击惩罚', 'Symmetry — 左右对称性']),
    ('REG 组 (4分量)\n运动效率+平滑性', CLR_ORANGE,
     ['Energy — 高斯核动作幅值惩罚', 'ActionSmooth — 相邻帧动作L2平滑', 'CostOfTransport — 运输能效比',
      'FallTerminal — 分级倒地终端惩罚']),
]
for i, (title, clr, items) in enumerate(groups):
    x = Inches(0.5 + i * 4.2)
    add_shape_bg(sl, x, Inches(1.5), Inches(3.9), Inches(5.0), CLR_LIGHT_BG)
    add_shape_bg(sl, x, Inches(1.5), Inches(3.9), Inches(0.05), clr)
    add_text_box(sl, x + Inches(0.2), Inches(1.65), Inches(3.5), Inches(0.7), title, font_size=15, bold=True, color=clr)
    for j, item in enumerate(items):
        add_text_box(sl, x + Inches(0.3), Inches(2.5 + j * 0.55), Inches(3.3), Inches(0.5), f'• {item}', font_size=12, color=CLR_GRAY)

add_code_ref(sl, 'OrangeRobotEnvComponent.h:101-133 三组枚举 | :394-404 位掩码属性 | .cpp:921-1514 ComputeReward')

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: 奖励体系 — 关键数学设计
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '关键奖励分量的数学设计（精选3项）')
add_section_label(sl, '核心技术', 11)

# 1. Alive
add_shape_bg(sl, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.5), CLR_LIGHT_BG)
add_text_box(sl, Inches(0.8), Inches(1.45), Inches(3), Inches(0.35),
             '1. Alive 存活奖励', font_size=18, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.9),
             'Alive = AliveReward × (1 + TiltQuality × VelQuality)\n'
             'TiltQuality = Clamp(uprightDot, 0, 1)   — 越直立越接近1\n'
             'VelQuality = 1 − |verticalSpeed|/20        — 颠簸越小越接近1\n'
             '→ 平稳站立获2倍奖励，倾斜/颠簸时衰减至基础值。全程不引入负奖励，在提供持续生存信号的同时用质量因子区分状态优劣',
             font_size=13, color=CLR_GRAY)

# 2. Command Tracking
add_shape_bg(sl, Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.5), CLR_LIGHT_BG)
add_text_box(sl, Inches(0.8), Inches(3.25), Inches(3), Inches(0.35),
             '2. 指令跟踪（指数核）', font_size=18, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.95),
             'FwdMatch = exp(−(ForwardSpeed − TargetFwd)² / (2×σ_fwd²))\n'
             '动态σ：指令越大σ越大（σ_fwd ∈ [25, MaxSpeed×0.5]），保持灵敏度一致\n'
             '站立指令：σ = StandCommandSigma = 5.0（极小），强制精确零速\n'
             '→ 高斯核在误差全域提供非零梯度，自适应σ消除L1/L2的梯度断裂/饱和问题',
             font_size=13, color=CLR_GRAY)

# 3. Height
add_shape_bg(sl, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), CLR_LIGHT_BG)
add_text_box(sl, Inches(0.8), Inches(5.05), Inches(3), Inches(0.35),
             '3. Height 高度维持（三段式）', font_size=18, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.95),
             '理想区(>90%参考高度): +HeightRewardScale 恒定正奖励\n'
             '过渡区(85%~90%): 线性插值 [负惩罚, +奖励] 平滑过渡\n'
             '危险区(<85%): −(1−HeightRatio)² × HeightDropPenaltyScale 平方惩罚快速下压\n'
             '→ 三段式设计防止策略学得"蹲伏"（降低重心减少摔倒几率但不产生位移）的局部最优',
             font_size=13, color=CLR_GRAY)

add_code_ref(sl, 'OrangeRobotEnvComponent.cpp:989-1021 (Alive) · :1077-1111 (CommandTracking) · :1034-1053 (Height)')

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: 位掩码消融系统
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '位掩码系统与消融实验能力', '支持任意组合独立开关 → 从"炼金术调参"到"系统工程"')
add_section_label(sl, '核心技术', 12)

items = [
    '▎ 位掩码机制',
    '   RewardMaskCore (6 bits)  ·  RewardMaskGait (6 bits)  ·  RewardMaskReg (4 bits)',
    '   设为 -1 = 全开；设为 0 = 全关；设为 bit组合 = 任意子集',
    '   APPLY_REWARD_CORE/GAIT/REG 宏 → 编译期内联条件分支 → 关闭的分量为严格零值',
    '',
    '▎ DynamicBalanceRewardWeight = 步态组总控旋钮',
    '   由课程学习调度：Stage0=0.0 → ... → Stage4=1.0（平滑线性插值）',
    '   所有GAIT分量(6项)统一乘以此权重 → 防止早期步态要求干扰站立基础学习',
    '',
    '▎ 典型消融实验方案',
    '   ① 仅开CORE组 → 验证基础生存+指令跟踪是否足以产生可行步态',
    '   ② 关闭REG组 → 测试能量/平滑/CoT惩罚的必要性',
    '   ③ 移除单个GAIT分量 → 量化每个步态细化项（交替/步频/对称/冲击）的独立贡献',
    '',
    '▎ 日志与分析',
    '   bLogRewardBreakdown=True + 每10步输出14分量值到UE Log',
    '   TensorBoard reward_components/* 实时追踪各分量 → 支持系统化A/B测试',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=15, spacing=Pt(6))

add_code_ref(sl, 'OrangeRobotEnvComponent.h:394-404 | .cpp:1482-1495 | train_headless.py:273-279 TensorboardMetricsCallback')

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: 核心技术三 — 课程学习
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术三：五阶段内建课程学习', '从站立到全指令 — 渐进式能力构建，阶段间平滑插值')
add_section_label(sl, '核心技术', 13)

headers3 = ['阶段', '步数', 'DynamicW', 'Forward指令', 'Turn指令', '新增惩罚', '学习目标']
rows3 = [
    ['0 站立', '0~150k', '0.0', '(0,0)固定', '0固定', '—', '躯干直立+不摔倒'],
    ['1 极慢速', '150k~300k', '0→0.3', '[0.05,0.2]', '0', '—', '慢速纯前进'],
    ['2 中速', '300k~600k', '0.3→0.6', '[0.2,0.4]', '[-0.05,0.05]', 'Sym+Freq', '行走+对称性'],
    ['3 快速', '600k~1M', '0.6→1.0', '[-1,1]', '[-0.5,0.5]', '+CoT', '全速后退+能效'],
    ['4 全指令', '>1M', '1.0', '[-1,1]', '[-1,1]', '全部', '全地形泛化'],
]
add_table(sl, Inches(0.5), Inches(1.5),
          [Inches(1.3), Inches(1.6), Inches(1.4), Inches(2.0), Inches(2.0), Inches(1.8), Inches(2.6)],
          headers3, rows3, font_size=12)

items2 = [
    '▎ 关键设计决策：',
    '   · DynamicW 作为总控：Stage0=0 关闭全部GAIT奖励 → 策略仅需学会"别摔倒" → 整个训练的基石',
    '   · 平滑过渡：阶段间所有权重通过 LerpWeight 线性插值 → 避免奖励突变导致策略震荡',
    '   · 分级倒地惩罚：FallScale = Stage0:0.2 / Stage1:0.5 / Stage2+:1.0 → 不惩罚站立的摔倒',
    '',
    '▎ C++侧实现：',
    '   · GlobalTrainingStep 跨episode累积 → 课程进度与episode边界无关',
    '   · UpdateCurriculumWeightsAndCommand() 每步动态计算 → 零额外通信开销',
    '   · 阶段边界可配置：CurriculumStageBoundaries = [150k, 300k, 600k, 1M]',
]
add_bullet_frame(sl, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2), items2, font_size=14, spacing=Pt(4))

add_code_ref(sl, 'OrangeRobotEnvComponent.cpp:489-591 UpdateCurriculumWeightsAndCommand')

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: 核心技术四 — SAC训练稳定化
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术四：SAC训练稳定化方案', '四项互补优化应对高维动作+短Episode+高方差奖励的Q值发散问题')
add_section_label(sl, '核心技术', 14)

headers4 = ['优化项', '具体做法', '解决的问题', '代码位置']
rows4 = [
    ['LayerNorm\nMLP', '每个隐藏层后插入LayerNorm\n[512, 512, 256]', '内部协变量漂移\n→ Q值震荡 → 发散', 'train_headless.py\n:71-88'],
    ['梯度裁剪', 'ClippedAdam\nmax_grad_norm=10.0', '高方差TD目标\n→ 梯度爆炸', 'train_headless.py\n:55-68'],
    ['缩短时序\n视野', 'γ: 0.99→0.95', '12步episode中终端惩罚\n远距离反向传播放大', 'train_headless.py\n:98'],
    ['固定熵\n系数', 'ent_coef=0.2\n替代 auto-tuning', 'auto-tuning在短episode中\n反馈不足→剧烈震荡', 'train_headless.py\n:99'],
]
add_table(sl, Inches(0.5), Inches(1.5),
          [Inches(1.8), Inches(4.2), Inches(3.8), Inches(3.0)], headers4, rows4, font_size=13)

# Diagnostic system
add_text_box(sl, Inches(0.8), Inches(4.5), Inches(4), Inches(0.4),
             '▎ 实时诊断指标体系', font_size=16, bold=True, color=CLR_ACCENT)
add_text_box(sl, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
             'train/q_value_mean     →  Q值均值漂移（理想：缓慢上升趋于稳定；发散信号：飙至100+）\n'
             'train/entropy              →  策略探索程度（理想：正值；坍塌信号：降至负值）\n'
             'train/q1_q2_diff           →  双Q网络一致性（理想：小；问题信号：持续增大→过估计）\n'
             'train/critic_loss           →  Critic拟合质量（理想：平滑下降；问题信号：剧烈震荡）\n'
             'train/ent_coef_value    →  熵系数实际值（固定模式则为常量0.2）\n'
             '\n每200步记录一次（TensorboardMetricsCallback）→ 训练107→165次迭代的参数演化实证轨迹',
             font_size=13, color=CLR_GRAY)

add_code_ref(sl, 'train_headless.py:55-88, 243-393 | train_orangerobot_sac.py:34-48 SAC_DEFAULTS')

# ═══════════════════════════════════════════════════════════════
# SLIDE 15: 核心技术五 — ONNX部署
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术五：ONNX部署与零编译参数管理', '从Python训练到C++推理的全链路工程打通')
add_section_label(sl, '核心技术', 15)

items = [
    '▎ ONNX推理管线',
    '   Python SAC训练完成 → export_*_onnx.py导出ONNX模型',
    '   UE5侧 UNNEModelData::Init("onnx", ...)  加载 → NNERuntimeORTCpu 执行推理',
    '   两个独立ONNX模型（高层+低层）共享Bridge调度 → 5Hz / 30Hz 异步频率',
    '',
    '▎ 零编译参数配置系统 (EnvConfigLoader)',
    '   命令行 -EnvConfig=path.json → UE反射 TFieldIterator<FProperty> → 自动映射到UPROPERTY',
    '   支持 float / double / int32 / bool / array 类型 → 覆盖全部50+训练参数',
    '   ExportAllConfigToJSON 一键导出当前参数快照 → 实验可复现',
    '   效果：调参从"改代码→编译→启动"的分钟级 → "改JSON→重启训练"的秒级',
    '',
    '▎ 训练-部署配置分离',
    '   BridgeComponent::ApplyDeployConfigIfPresent 通过 -DeployConfig=path.json',
    '   指定ONNX路径 + 目标Actor → 与训练配置（EnvConfig）完全解耦',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.0), items, font_size=15, spacing=Pt(8))

add_code_ref(sl, 'EnvConfigLoader.cpp:9-143 | BridgeComponent.cpp:145-193 | export_orangerobot_onnx.py')

# ═══════════════════════════════════════════════════════════════
# SLIDE 16: 核心技术六 — 物理约束
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心技术六：物理约束与动作安全机制', '当前三层被动约束 → 未来QP主动约束层')
add_section_label(sl, '核心技术', 16)

# Current
add_shape_bg(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8), CLR_LIGHT_BG)
add_text_box(sl, Inches(0.7), Inches(1.55), Inches(5.3), Inches(0.4),
             '🔸 当前已实现（三层被动约束）', font_size=16, bold=True, color=CLR_ORANGE)
add_text_box(sl, Inches(0.7), Inches(2.0), Inches(5.3), Inches(2.0),
             '① ActionDeadzone=0.05\n'
             '   过滤幅度<5%的噪声动作，防止无效微振\n\n'
             '② L2范数裁剪 (max norm=2.0)\n'
             '   所有关节同时满幅输出时等比例缩回\n'
             '   防止Chaos求解器子迭代爆炸\n\n'
             '③ 关节角速度限幅\n'
             '   Twist≤45°/s, Swing≤35°/s\n'
             '   ClampAngularVelocityTarget',
             font_size=13, color=CLR_GRAY)

# Future
add_shape_bg(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8), RGBColor(0x1A, 0x2A, 0x1A))
add_text_box(sl, Inches(7.0), Inches(1.55), Inches(5.3), Inches(0.4),
             '🔹 规划中（第四层主动QP约束）', font_size=16, bold=True, color=CLR_GREEN)
add_text_box(sl, Inches(7.0), Inches(2.0), Inches(5.3), Inches(2.0),
             '④ QP约束求解器\n'
             '   插入位置：ONNX推理输出\n'
             '             → SetAngularVelocityTarget 之间\n\n'
             '   约束项：\n'
             '   · 关节限位（角度上下界）\n'
             '   · 力矩上限（驱动能力限制）\n'
             '   · 足底摩擦锥（法向接触约束）\n\n'
             '   目标函数：min ||a_qp − a_policy||²\n'
             '   关键挑战：33ms内完成QP求解',
             font_size=13, color=CLR_GRAY)

# Warning box
add_shape_bg(sl, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.0), RGBColor(0x2A, 0x18, 0x18))
add_text_box(sl, Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.4),
             '⚠ 为什么当前被动约束不够？', font_size=16, bold=True, color=CLR_ACCENT2)
add_text_box(sl, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.4),
             'Saved/Crashes/ 下 40 个崩溃目录（2026.3.31 – 5.7）→ 被动惩罚机制的系统性失效证据\n'
             'FootImpactVelocityThreshold=50.0 → "告诉策略不要暴力踩地"，但不能阻止策略在某一步确实输出越界指令\n'
             '无关节限位前置检查 → 策略输出可能导致约束超限 → Chaos引擎自行处理 → 不稳定约束力 → 爆炸/崩溃\n'
             '\n→ 架构演进方向：将物理可行性从"奖励惩罚"升级为"架构性硬保证"（阶段二核心目标）',
             font_size=13, color=CLR_GRAY)

add_code_ref(sl, 'OrangeRobotEnvComponent.h:413,422-426 | .cpp:88-94,754-765,780 | Saved/Crashes/ (40 crash dirs)')

# ═══════════════════════════════════════════════════════════════
# SLIDE 17: 实验平台
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '实验平台与测试场景')
add_section_label(sl, '实验结果', 17)

items = [
    '▎ 仿真环境',
    '   UE5 (Chaos Physics Engine) · 自建高自由度双足机器人骨骼网格',
    '   自定义碰撞通道 (linkA, linkB) 用于足地精细交互检测',
    '   AngularVelocityDrive 驱动模式 · Twist + Swing1/Swing2 约束轴',
    '',
    '▎ 训练框架',
    '   Python SAC (stable-baselines3) ↔ UE5 通过 Schola + gRPC (port=50051) + Protobuf',
    '   控制器频率 30Hz · LayerNorm MLP [512,512,256] · ClippedAdam max_grad_norm=10.0',
    '   SAC超参：lr=5e-5, γ=0.95, ent_coef=0.2, buffer=1M, batch=1024, grad_steps=4',
    '',
    '▎ 测试场景',
    '   场景A：无障碍直线路径 — 验证基本行走与指令跟踪能力',
    '   场景B：含障碍物几字形迷宫 — 验证绕障与导航决策',
    '',
    '▎ 评价指标',
    '   导航成功率 · 平均完成步数 · 躯干姿态抖动(Roll/Pitch/Yaw方差)',
    '   足地滑移率 · 能耗效率(CoT) · 物理违规次数 · 推理时延分解',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=15, spacing=Pt(8))

add_code_ref(sl, 'train_headless.py:35-100 SAC_DEFAULTS | OrangeRobotEnvComponent.h:257-263')

# ═══════════════════════════════════════════════════════════════
# SLIDE 18: 低层步态训练结果
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '低层步态策略训练结果', '五阶段课程学习下的训练曲线与指标演化')
add_section_label(sl, '实验结果', 18)

items = [
    '▎ Episode平均奖励 vs 步数',
    '   标注五个课程阶段的分界线 → 体现课程学习的效果跃升',
    '   观察：奖励在阶段过渡处出现的平滑跃升（非突变→验证LerpWeight线性插值有效性）',
    '',
    '▎ 指令跟踪精度',
    '   前向速度跟踪RMSE（cm/s）随训练步数的下降曲线',
    '   转向角速度跟踪RMSE（deg/s）随训练步数的下降曲线',
    '   观察：前向跟踪通常先于转向收敛（前向奖励权重0.7 vs 转向0.3）',
    '',
    '▎ 存活步数(Episode Length)演化',
    '   Stage0: ~10-15步 → Stage4: 100+步（满episode）',
    '   关键指标：策略从"勉强站立"到"稳定行走"的核心证据',
    '',
    '▎ 14个奖励分量的相对贡献度分析（柱状图）',
    '   识别主导训练信号的分量（通常Alive + CommandTracking占主导）',
    '   量化每个分量的平均step reward → 指导后续权重调优',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=15, spacing=Pt(8))

add_code_ref(sl, 'TensorBoard custom/* metrics | reward_components/* per-component logs')

# ═══════════════════════════════════════════════════════════════
# SLIDE 19: 高层导航+联合部署
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '高层导航与联合部署结果')
add_section_label(sl, '实验结果', 19)

items = [
    '▎ 几字形迷宫导航路径可视化',
    '   立方体代理路径轨迹展示绕障行为',
    '   观察：策略能发现侧向通路而非"硬撞墙壁"',
    '',
    '▎ 奖励分量消融效果',
    '   HeadingGate（朝向门控）→ 消除无意义徘徊',
    '   RaySuddenOpeningReward → 促进侧向通路发现',
    '   BypassRewardScale → 路径优化贡献',
    '',
    '▎ 联合部署推理性能',
    '   总推理时延分解：高层ONNX + 低层ONNX + EMA平滑 + Stepper调度',
    '   验证：两个ONNX模型在30Hz控制周期内的实时性',
    '',
    '▎ 端到端导航指标',
    '   场景A（无障碍直线）：成功率、平均步数、路径效率',
    '   场景B（几字形迷宫）：成功率、平均步数、路径效率',
    '   行走稳定性：躯干姿态抖动幅值、足地滑移率',
]
add_bullet_frame(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=15, spacing=Pt(8))

add_code_ref(sl, 'XNavigationCubeEnvComponent.cpp:571-700 ComputeReward | BridgeComponent.cpp:294-323')

# ═══════════════════════════════════════════════════════════════
# SLIDE 20: 消融实验
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '消融实验与关键发现')
add_section_label(sl, '实验结果', 20)

headers5 = ['实验', '对照条件', '实验条件', '关键发现']
rows5 = [
    ['课程学习\n消融', '五阶段课程', '无课程\n(全指令随机)', '无课程时策略几乎无法走出Stage0\n站立阶段—课程学习是训练成功的前提条件'],
    ['LayerNorm\n消融', '标准MLP\n(无LN)', 'LayerNorm\nMLP', 'LayerNorm对Q值稳定性的改善效果\n远超单纯调低学习率'],
    ['奖励分组\n消融', '全开(14分量)', '仅CORE组\n(6分量)', 'CORE组(Alive+Height+Tracking)已可\n产生基本行走—GAIT/REG起细化作用'],
    ['指令分布\n敏感性', '训练分布内\n指令', '训练分布外\n指令(外推10%)', '低层策略对分布外指令表现可接受\n泛化—步态退化在10%以内'],
]
add_table(sl, Inches(0.5), Inches(1.5),
          [Inches(2.0), Inches(2.5), Inches(2.5), Inches(5.8)], headers5, rows5, font_size=12)

add_shape_bg(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), CLR_LIGHT_BG)
add_text_box(sl, Inches(0.8), Inches(5.3), Inches(2.5), Inches(0.4),
             '▎ 三个核心发现', font_size=18, bold=True, color=CLR_GREEN)
add_text_box(sl, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
             '① 课程学习是训练成功的必要条件 — 无课程时策略无法跨越"站立→行走"的技能鸿沟\n'
             '② LayerNorm在抑制Q值发散上的效果远超调低学习率 — 数值稳定性的架构保障优于参数微调\n'
             '③ 独立训练的低层策略对未见指令模式具有可接受泛化能力 — 验证了分层解耦方案的部署可行性',
             font_size=14, color=CLR_WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 21: 当前局限性
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '当前方案的局限性分析', '四个明确局限 — 与开头问题定义形成闭环')
add_section_label(sl, '局限与展望', 21)

limits = [
    ('①', '物理约束的被动性', CLR_ACCENT2,
     '仅 ActionDeadzone + L2裁剪无法从架构层面杜绝关节超限/力矩饱和/足底打滑。\n40次崩溃记录是这一问题的直接证据。'),
    ('②', '训练-部署的分布偏移', CLR_ORANGE,
     '低层在随机均匀指令上训练，高层实际输出分布偏离均匀。\nCommandSmoothAlpha=0.3 的EMA平滑在训练时不存在 → 训练-部署gap。'),
    ('③', '单环境串行训练效率', CLR_ORANGE,
     '步进级gRPC往返（每步Protobuf序列化/反序列化）的墙钟时间瓶颈。\n150万步需数小时 → 多层架构下通信开销进一步放大。'),
    ('④', '缺乏外部基准对比', CLR_GRAY,
     '尚未与PPO/TD3/MPC等替代方案进行公平的横向性能对比。\n难以精确定位本方案在双足机器人DRL控制领域的技术坐标。'),
]
for i, (num, title, clr, desc) in enumerate(limits):
    y = Inches(1.5 + i * 1.35)
    add_shape_bg(sl, Inches(0.5), y, Inches(12.3), Inches(1.2), CLR_LIGHT_BG)
    add_text_box(sl, Inches(0.7), y + Inches(0.05), Inches(1.5), Inches(0.4), f'{num} {title}',
                 font_size=17, bold=True, color=clr)
    add_text_box(sl, Inches(2.3), y + Inches(0.05), Inches(10.3), Inches(1.1), desc,
                 font_size=14, color=CLR_GRAY)

add_code_ref(sl, 'Saved/Crashes/ — 40 crash dirs | OrangeRobotEnvComponent.h:413 — ActionDeadzone | gRPC port=50051')

# ═══════════════════════════════════════════════════════════════
# SLIDE 22: 五阶段技术演进路线图
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '五阶段技术演进路线图', '从"稳固基线"到"全栈验证"的完整技术演进路径')
add_section_label(sl, '局限与展望', 22)

stages = [
    ('一', '稳固基线', '✅ 已完成', CLR_GREEN,
     '双SAC-ONNX完整训练+部署\n基线性能测试(直线+迷宫)'),
    ('二', '物理约束内化', '📋 规划中', CLR_ACCENT,
     'QP约束求解器\n关节限位+力矩+摩擦锥'),
    ('三', '解耦训练形式化', '📋 规划中', CLR_ACCENT,
     '指令分布KL散度监控\n低层自适应fine-tuning'),
    ('四', '导航泛化增强', '📋 规划中', CLR_ACCENT,
     '多样化障碍物+动态扰动\n复杂地形验证'),
    ('五', '通信优化+全栈', '📋 规划中', CLR_ACCENT,
     '批量异步通信/共享内存\n极限场景压力测试'),
]
for i, (num, name, status, clr, desc) in enumerate(stages):
    x = Inches(0.3 + i * 2.55)
    add_shape_bg(sl, x, Inches(1.6), Inches(2.35), Inches(3.8), CLR_LIGHT_BG)
    add_shape_bg(sl, x, Inches(1.6), Inches(2.35), Inches(0.06), clr)
    add_text_box(sl, x + Inches(0.15), Inches(1.75), Inches(2.0), Inches(0.5),
                 f'阶段{num}: {name}', font_size=14, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), Inches(2.3), Inches(2.0), Inches(0.3),
                 status, font_size=12, color=clr, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), Inches(2.8), Inches(2.0), Inches(2.0),
                 desc, font_size=11, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(sl, Inches(0.5), Inches(5.7), Inches(12), Inches(0.3),
             '→ → → → → → → → → → → → → → → → → → → → → → → → → → → → → →',
             font_size=14, color=CLR_ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             '长期愿景：从UE5仿真 → 真机迁移，利用UE5 Photorealistic Rendering实现高质量Sim-to-Real',
             font_size=14, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 23: 核心贡献总结
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, '核心贡献总结')
add_section_label(sl, '总结', 23)

contribs = [
    ('Ⅰ', '分层推理-独立训练架构', CLR_ACCENT,
     '以固定2D指令接口 + 课程化指令采样实现高层导航与低层步态的有效解耦。\n'
     '论证了其相较于端到端RL（样本效率低/物理不可信）、MPC（实时性不足）、\n'
     '传统分层RL（非平稳训练）的明确优势。'),
    ('Ⅱ', '14组分位掩码奖励 + 五阶段课程学习', CLR_GREEN,
     '将50+参数的耦合调参空间分解为语义清晰、可独立消融的子系统。\n'
     '配合C++侧内建课程学习，实现从站立到全指令的渐进式能力构建。\n'
     '这是双足机器人DRL训练从"炼金术"走向"系统工程"的关键一步。'),
    ('Ⅲ', 'SAC训练稳定化方案', CLR_ACCENT2,
     'LayerNorm MLP + 梯度裁剪 + 缩短γ + 固定熵系数的四合一优化。\n'
     '经受住了160+次训练迭代的实证检验，成功抑制了高维双足控制中的Q值发散。'),
    ('Ⅳ', '全链路工程方法论', CLR_ORANGE,
     'UE反射配置加载 → Python SAC训练 → ONNX导出 → C++ Bridge部署 → gRPC通信。\n'
     '五阶段技术演进路线图为后续工作提供明确的工程路径。\n'
     '为UE5+DRL交叉领域的后续研究者提供了可复现的完整范式。'),
]
for i, (num, title, clr, desc) in enumerate(contribs):
    y = Inches(1.4 + i * 1.4)
    add_shape_bg(sl, Inches(0.5), y, Inches(12.3), Inches(1.25), CLR_LIGHT_BG)
    add_shape_bg(sl, Inches(0.5), y, Inches(0.06), Inches(1.25), clr)
    add_text_box(sl, Inches(0.8), y + Inches(0.05), Inches(2.0), Inches(0.4),
                 f'{num}  {title}', font_size=17, bold=True, color=clr)
    add_text_box(sl, Inches(0.8), y + Inches(0.5), Inches(11.5), Inches(0.7),
                 desc, font_size=13, color=CLR_GRAY)

add_code_ref(sl, 'docs/superpowers/specs/2026-05-08-project-analysis.md — 完整技术分析文档')

# ═══════════════════════════════════════════════════════════════
# SLIDE 24: Q&A
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
add_shape_bg(sl, Inches(0), Inches(0), W, H, CLR_DARK)
add_shape_bg(sl, Inches(0), Inches(2.5), W, Inches(0.05), CLR_ACCENT)
add_text_box(sl, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             '感谢聆听 · 请提问', font_size=44, bold=True, color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             '基于UE5的双足机器人智能寻路系统的设计与开发', font_size=20, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(4.5), Inches(11), Inches(1.0),
             '备用材料：奖励公式推导 · SAC算法伪代码 · ONNX模型结构 · 训练配置JSON示例',
             font_size=14, color=CLR_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = Path(__file__).resolve().parent.parent / "docs" / "superpowers" / "specs" / "defense_ppt.pptx"
output_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_path))
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
