"""Insert arch diagram slide into defense PPT at position 7."""
from pptx import Presentation
from pathlib import Path
from copy import deepcopy
from lxml import etree

base = Path(__file__).resolve().parent.parent
defense_path = base / "docs" / "superpowers" / "specs" / "defense_ppt.pptx"
arch_path = base / "docs" / "superpowers" / "specs" / "arch_diagram.pptx"

# Load both
prs_defense = Presentation(str(defense_path))
prs_arch = Presentation(str(arch_path))

# Get arch slide (only slide)
arch_slide = prs_arch.slides[0]

# Insert arch slide's XML at position 6 (after slide 6, making it slide 7)
# We need to manipulate the XML directly
slide_to_insert = arch_slide

# Get the slide XML
arch_slide_xml = slide_to_insert._element

# Get the sldIdLst
presentation_xml = prs_defense.part._element
sldIdLst = presentation_xml.find('{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')

# Get slide 6's sldId element (0-indexed: index 5)
sld_ids = list(sldIdLst)
target_sld_id = sld_ids[5]  # slide 6 (0-indexed)

# Create new slide ID
max_id = max(int(s.get('id')) for s in sld_ids)
new_sld_id = str(max_id + 1)

# Build new sldId element
nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
         'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
new_sld = etree.SubElement(sldIdLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
new_sld.set('id', new_sld_id)
# We need a r:id - get the next available rId
# First, move new_sld to correct position (after target_sld_id)
sldIdLst.remove(new_sld)
target_sld_id.addnext(new_sld)

# Add slide relationship and part
slide_part_name = f'/ppt/slides/slide{len(prs_defense.slides) + 1}.xml'
# Actually, let's use a different approach - add the slide via the standard API
# but reorder it

# Remove the XML hack and re-add the slide properly
sldIdLst.remove(new_sld)

# Instead, let's add via API then reorder
prs_defense.slides.add_slide(prs_defense.slide_layouts[6])  # blank placeholder
# We'll replace its content with arch slide content

# Actually the cleanest approach: add the arch slide to defense, reorder
# Let me use a simpler approach - clone the arch_slide's content into the defense

# Delete the temp blank slide we just added
# This is getting complex with python-pptx XML manipulation.
# Simpler approach: rebuild defense PPT with arch slide inserted.

# Let me just use a direct approach
prs = Presentation(str(defense_path))
prs_arch2 = Presentation(str(arch_path))

slide_count = len(prs.slides)
# Add arch slide content to a new blank slide
blank_layout = prs.slide_layouts[6]  # blank
new_slide = prs.slides.add_slide(blank_layout)

# Copy all shapes from arch slide to new slide
arch_sl = prs_arch2.slides[0]
for shape in arch_sl.shapes:
    el = deepcopy(shape._element)
    new_slide.shapes._spTree.append(el)

# Now reorder: move the last slide to position 6 (after slide 5, making it slide 7)
slides_list = list(prs.slides._sldIdLst)
last_entry = slides_list[-1]
prs.slides._sldIdLst.remove(last_entry)
# Insert after index 5 (which is the 6th element, 0-indexed)
slides_list2 = list(prs.slides._sldIdLst)
target = slides_list2[5]
target.addnext(last_entry)

prs.save(str(defense_path))
print(f"Done! Slides: {len(prs.slides)}")
