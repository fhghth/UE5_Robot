"""
Generate a standalone system architecture diagram slide for the defense PPT.
Adds this slide as a replacement for slide 7 in the existing PPT.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pathlib import Path

# ── Colors ────────────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x0D, 0x1A)
C_PANEL  = RGBColor(0x16, 0x16, 0x2E)
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_ACCENT2= RGBColor(0xF7, 0x25, 0x5D)   # red
C_GREEN  = RGBColor(0x2E, 0xCC, 0x71)
C_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
C_WHITE  = RGBColor(0xF0, 0xF0, 0xF5)
C_GRAY   = RGBColor(0x90, 0x90, 0xA0)
C_LGRAY  = RGBColor(0x60, 0x60, 0x80)
C_HLBOX  = RGBColor(0x1E, 0x2A, 0x3E)
C_TRAIN  = RGBColor(0x1A, 0x28, 0x18)   # training panel bg

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Full dark background ──
bg = sl.background
bg.fill.solid()
bg.fill.fore_color.rgb = C_DARK

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════
def box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = border_width
    else:
        s.line.fill.background()
    return s

def txt(slide, left, top, width, height, text, size=11, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, name='Consolas'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb

def multi_txt(slide, left, top, width, height, lines, size=11, color=C_WHITE, line_spacing=Pt(4)):
    """lines is list of (text, size_override, color_override, bold) or just str"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
        else:
            p.text = line[0]
            p.font.size = Pt(line[1] if len(line) > 1 else size)
            p.font.color.rgb = line[2] if len(line) > 2 else color
            if len(line) > 3:
                p.font.bold = line[3]
        p.font.name = 'Microsoft YaHei'
        p.space_after = line_spacing
    return tb

def arrow_right(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5)):
    """Draw a right arrow connector"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # add arrowhead
    connector.line._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd'] = 'none'
    return connector

def arrow_down(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5)):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

# ═══════════════════════════════════════════════════════════
# TITLE BAR
# ═══════════════════════════════════════════════════════════
box(sl, Inches(0), Inches(0), W, Inches(0.7), C_PANEL)
box(sl, Inches(0.5), Inches(0.65), Inches(2.0), Inches(0.04), C_ACCENT)
txt(sl, Inches(0.8), Inches(0.1), Inches(11), Inches(0.5),
    '系统总体架构图', size=24, bold=True, color=C_WHITE)
txt(sl, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.3),
    '分层推理 · 独立训练 · 桥接部署', size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# LEFT SIDE: Training (Python)
# ═══════════════════════════════════════════════════════════
BOX_L = Inches(0.3)
BOX_W = Inches(3.8)

# Training panel background
box(sl, Inches(0.15), Inches(0.85), Inches(4.1), Inches(6.4), C_TRAIN, border_color=RGBColor(0x2A, 0x4A, 0x2A))
txt(sl, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.35),
    '🖥 训练端 (Python)', size=14, bold=True, color=C_GREEN)

# High-level training
box(sl, BOX_L, Inches(1.4), BOX_W, Inches(1.5), C_PANEL, border_color=C_ACCENT)
multi_txt(sl, Inches(0.5), Inches(1.45), Inches(3.4), Inches(1.4), [
    ('高层导航策略训练', 13, C_ACCENT, True),
    ('━━━━━━━━━━━━━━━━━━━', 9, C_GRAY),
    ('算法: SAC (stable-baselines3)', 10, C_GRAY),
    ('观测: 目标方向+距离+8射线 (17维)', 10, C_GRAY),
    ('动作: [Forward, Turn] 2维连续指令', 10, C_GRAY),
    ('环境: XNavigationCubeEnvComponent', 10, C_LGRAY),
])

# Low-level training
box(sl, BOX_L, Inches(3.1), BOX_W, Inches(1.8), C_PANEL, border_color=C_ORANGE)
multi_txt(sl, Inches(0.5), Inches(3.15), Inches(3.4), Inches(1.7), [
    ('低层步态策略训练', 13, C_ORANGE, True),
    ('━━━━━━━━━━━━━━━━━━━', 9, C_GRAY),
    ('算法: SAC · lr=5e-5 · γ=0.95', 10, C_GRAY),
    ('观测: 躯干+关节+足地+步态相位', 10, C_GRAY),
    ('      + 高层指令[Forward, Turn] (30-40维)', 10, C_GRAY),
    ('动作: 10-12维关节角速度目标', 10, C_GRAY),
    ('环境: OrangeRobotEnvComponent', 10, C_LGRAY),
])

# Training optimizations
box(sl, BOX_L, Inches(5.1), BOX_W, Inches(1.0), C_PANEL, border_color=C_GRAY)
multi_txt(sl, Inches(0.5), Inches(5.15), Inches(3.4), Inches(0.9), [
    ('训练优化', 12, C_GREEN, True),
    ('LayerNorm MLP [512,512,256]', 10, C_GRAY),
    ('ClippedAdam (max_grad_norm=10.0)', 10, C_GRAY),
    ('14组分奖励 · 位掩码消融 · 5阶段课程', 10, C_GRAY),
])

# Export
box(sl, BOX_L, Inches(6.3), BOX_W, Inches(0.7), C_PANEL, border_color=C_ACCENT2)
txt(sl, Inches(0.5), Inches(6.35), Inches(3.4), Inches(0.55),
    '📦 ONNX导出: export_*_onnx.py\n'
    '   高层.onnx + 低层.onnx',
    size=10, color=C_GRAY)

# ═══════════════════════════════════════════════════════════
# CENTER: gRPC Communication Bridge
# ═══════════════════════════════════════════════════════════
CX = Inches(4.5)
CW = Inches(4.3)

# gRPC box
box(sl, CX, Inches(2.2), CW, Inches(2.8), C_PANEL, border_color=C_ACCENT, border_width=Pt(2))
txt(sl, CX + Inches(0.2), Inches(2.3), Inches(3.9), Inches(0.35),
    '🔗 Schola gRPC 通信层 (Port=50051)', size=14, bold=True, color=C_ACCENT)
multi_txt(sl, CX + Inches(0.2), Inches(2.75), Inches(3.9), Inches(2.0), [
    ('训练时:', 12, C_GREEN, True),
    ('  Python SAC ⇄ UE5 步进级通信', 10, C_GRAY),
    ('  Protobuf 序列化/反序列化', 10, C_GRAY),
    ('  每步：观测→gRPC→推理→gRPC→动作', 10, C_GRAY),
    ('', 8, C_GRAY),
    ('推理/部署时:', 12, C_ORANGE, True),
    ('  旁路gRPC → ONNX直接在C++侧推理', 10, C_GRAY),
    ('  UNNEPolicy + NNERuntimeORTCpu', 10, C_GRAY),
    ('  消除网络延迟 → 满足30Hz实时性', 10, C_GRAY),
])

# ═══════════════════════════════════════════════════════════
# RIGHT SIDE: UE5 Deployment (C++)
# ═══════════════════════════════════════════════════════════
RX = Inches(9.1)
RW = Inches(4.0)

# UE5 panel background
box(sl, Inches(8.95), Inches(0.85), Inches(4.2), Inches(6.4), RGBColor(0x18, 0x18, 0x2A), border_color=RGBColor(0x3A, 0x3A, 0x5A))
txt(sl, Inches(9.3), Inches(0.9), Inches(3.5), Inches(0.35),
    '🎮 部署端 (UE5 C++)', size=14, bold=True, color=C_ACCENT)

# Bridge Component
box(sl, RX, Inches(1.4), RW, Inches(1.5), C_PANEL, border_color=C_ACCENT, border_width=Pt(2))
multi_txt(sl, Inches(9.3), Inches(1.45), Inches(3.5), Inches(1.4), [
    ('⚙ ABridgeComponent (核心调度器)', 13, C_ACCENT, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━', 9, C_GRAY),
    ('· LoadModelDataFromDisk 加载双ONNX', 10, C_GRAY),
    ('· InitPolicies → 高层UNNEPolicy + 低层UNNEPolicy', 10, C_GRAY),
    ('· InitSteppers → USimpleStepper(高层)', 10, C_GRAY),
    ('               + UPipelinedStepper(低层)', 10, C_GRAY),
    ('· 调度逻辑: Tick() 每帧调用', 10, C_GRAY),
])

# High-level inference
box(sl, RX, Inches(3.1), RW, Inches(0.9), C_PANEL, border_color=C_ACCENT)
multi_txt(sl, Inches(9.3), Inches(3.15), Inches(3.5), Inches(0.8), [
    ('🔹 StepHighLevelInference (5Hz)', 12, C_ACCENT, True),
    (' HighInferenceInterval = 0.2s', 10, C_GRAY),
    (' NavigationStepper→Step() → RawHighLevelCommand', 10, C_GRAY),
    (' → EMA平滑 (α=0.3) → SmoothedCommand', 10, C_GRAY),
])

# Low-level inference
box(sl, RX, Inches(4.2), RW, Inches(1.0), C_PANEL, border_color=C_ORANGE)
multi_txt(sl, Inches(9.3), Inches(4.25), Inches(3.5), Inches(0.9), [
    ('🔸 StepLowLevelInference (30Hz)', 12, C_ORANGE, True),
    (' ControlStepper→Step() 每帧执行', 10, C_GRAY),
    (' 观测: 本体感知(30-40维) + HighLevelCommand(2维)', 10, C_GRAY),
    (' → 输出10-12维关节角速度目标', 10, C_GRAY),
])

# Physical Safety Layer
box(sl, RX, Inches(5.4), RW, Inches(0.9), RGBColor(0x28, 0x18, 0x18), border_color=C_ACCENT2)
multi_txt(sl, Inches(9.3), Inches(5.45), Inches(3.5), Inches(0.8), [
    ('🛡 动作安全层', 12, C_ACCENT2, True),
    (' ① ActionDeadzone=0.05  过滤微小噪声', 10, C_GRAY),
    (' ② L2范数裁剪(max=2.0) 防求解器爆炸', 10, C_GRAY),
    (' ③ 角速度限幅 Twist≤45° Swing≤35°', 10, C_GRAY),
])

# Chaos Engine
box(sl, RX, Inches(6.5), RW, Inches(0.55), C_ACCENT2)
txt(sl, Inches(9.3), Inches(6.52), Inches(3.5), Inches(0.5),
    '⚡ Chaos Physics Engine (30Hz)',
    size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# ARROWS & FLOW ANNOTATIONS
# ═══════════════════════════════════════════════════════════

# Training → gRPC arrows
txt(sl, Inches(4.1), Inches(2.0), Inches(0.6), Inches(0.3), '▸', size=18, color=C_GREEN, align=PP_ALIGN.CENTER)
txt(sl, Inches(4.1), Inches(3.8), Inches(0.6), Inches(0.3), '▸', size=18, color=C_ORANGE, align=PP_ALIGN.CENTER)

# gRPC → UE5 arrows
txt(sl, Inches(8.6), Inches(2.0), Inches(0.6), Inches(0.3), '▸', size=18, color=C_GREEN, align=PP_ALIGN.CENTER)
txt(sl, Inches(8.6), Inches(3.8), Inches(0.6), Inches(0.3), '▸', size=18, color=C_ORANGE, align=PP_ALIGN.CENTER)

# ONNX deployment arrow
txt(sl, Inches(6.2), Inches(6.6), Inches(1.5), Inches(0.3), '── ONNX直接加载 ──▶', size=9, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Internal Bridge flow arrows (right side vertical)
txt(sl, Inches(10.8), Inches(2.9), Inches(0.5), Inches(0.3), '▼', size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, Inches(10.8), Inches(4.0), Inches(0.5), Inches(0.3), '▼', size=14, color=C_ORANGE, align=PP_ALIGN.CENTER)
txt(sl, Inches(10.8), Inches(5.2), Inches(0.5), Inches(0.3), '▼', size=14, color=C_ACCENT2, align=PP_ALIGN.CENTER)
txt(sl, Inches(10.8), Inches(6.2), Inches(0.5), Inches(0.3), '▼', size=14, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# FOOTER: Code references
# ═══════════════════════════════════════════════════════════
txt(sl, Inches(0.5), Inches(7.2), Inches(12), Inches(0.25),
    'BridgeComponent.h/cpp · OrangeRobotEnvComponent.h/cpp · XNavigationCubeEnvComponent.h/cpp · train_headless.py · EnvConfigLoader.cpp',
    size=8, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output_path = Path(__file__).resolve().parent.parent / "docs" / "superpowers" / "specs" / "arch_diagram.pptx"
prs.save(str(output_path))
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
