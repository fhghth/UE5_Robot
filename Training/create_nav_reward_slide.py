"""
Generate a single PPT slide: Navigation Reward Function Design
Insert into defense PPT at position 10 (after low-level reward slides, before curriculum)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from copy import deepcopy

# ── Colors ────────────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x0D, 0x1A)
C_PANEL  = RGBColor(0x16, 0x16, 0x2E)
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2= RGBColor(0xF7, 0x25, 0x5D)
C_GREEN  = RGBColor(0x2E, 0xCC, 0x71)
C_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
C_PURPLE = RGBColor(0xA0, 0x60, 0xF0)
C_WHITE  = RGBColor(0xF0, 0xF0, 0xF5)
C_GRAY   = RGBColor(0x90, 0x90, 0xA0)
C_LGRAY  = RGBColor(0x60, 0x60, 0x80)
C_BG2    = RGBColor(0x22, 0x22, 0x3E)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
sl = prs.slides.add_slide(prs.slide_layouts[6])

bg = sl.background
bg.fill.solid()
bg.fill.fore_color.rgb = C_DARK

def box(slide, left, top, width, height, fill_color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, left, top, width, height, text, size=11, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tb

def mtext(slide, left, top, width, height, lines, base_size=10, base_color=C_GRAY):
    """lines: list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[0]
        p.font.size = Pt(line[1] if len(line) > 1 else base_size)
        p.font.color.rgb = line[2] if len(line) > 2 else base_color
        if len(line) > 3: p.font.bold = line[3]
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(3)
    return tb

# ═══════════════════════════════════════════════════════════
# TITLE BAR
# ═══════════════════════════════════════════════════════════
box(sl, Inches(0), Inches(0), W, Inches(0.65), C_PANEL)
box(sl, Inches(0.5), Inches(0.6), Inches(2.0), Inches(0.04), C_PURPLE)
txt(sl, Inches(0.8), Inches(0.08), Inches(11), Inches(0.5),
    '导航奖励函数设计：15分量 + 5组架构', size=22, bold=True, color=C_WHITE)
txt(sl, Inches(10.5), Inches(0.12), Inches(2.5), Inches(0.3),
    'XNavigationCubeEnvComponent::ComputeReward', size=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# LEFT: Architecture overview + key design philosophy
# ═══════════════════════════════════════════════════════════
LX = Inches(0.35)
LW = Inches(3.8)

box(sl, LX, Inches(0.85), LW, Inches(6.4), C_BG2, border_color=C_PURPLE)

# Design philosophy
mtext(sl, LX + Inches(0.15), Inches(0.95), LW - Inches(0.3), Inches(1.6), [
    ('▎ 设计哲学', 14, C_PURPLE, True),
    ('━━━━━━━━━━━━━━━━━━', 9, C_LGRAY),
    ('导航奖励的核心矛盾:', 11, C_WHITE),
    ('  "尽快到达目标" vs "避开障碍物"', 10, C_GRAY),
    ('  → 单一距离奖励导致"撞墙冲刺"', 10, C_GRAY),
    ('', 6, C_GRAY),
    ('解法: 密集塑形奖励 + 方向门控', 11, C_GREEN),
    ('  · 距离缩短 × 朝向对齐门控', 10, C_GRAY),
    ('  · 多层级停滞检测(位置停滞+距离停滞)', 10, C_GRAY),
    ('  · 侧向通路发现奖励(射线突开+绕障机会)', 10, C_GRAY),
    ('  · 安全前瞻约束(盲目前冲惩罚)', 10, C_GRAY),
])

# Key parameters
mtext(sl, LX + Inches(0.15), Inches(2.65), LW - Inches(0.3), Inches(1.5), [
    ('▎ 关键参数', 14, C_PURPLE, True),
    ('━━━━━━━━━━━━━━━━━━', 9, C_LGRAY),
    ('MoveStepScale       = 20 cm', 10, C_GRAY),
    ('TurnStepDegrees      = 15°', 10, C_GRAY),
    ('ReachTargetDistance  = 50 cm', 10, C_GRAY),
    ('MaxObserveDistance   = 2000 cm', 10, C_GRAY),
    ('PerceptionHalfExtent = 30 cm (盒体Sweep)', 10, C_GRAY),
    ('NumRays              = 8 (等角45°)', 10, C_GRAY),
    ('MaxConsecutiveStuckSteps = 25', 10, C_GRAY),
    ('HeadingGateExponent  = 2.0 (二次门控)', 10, C_GRAY),
])

# Observation space
mtext(sl, LX + Inches(0.15), Inches(4.3), LW - Inches(0.3), Inches(1.2), [
    ('▎ 观测空间 (17维)', 14, C_PURPLE, True),
    ('━━━━━━━━━━━━━━━━━━', 9, C_LGRAY),
    ('目标局部方向(3)  +  归一化距离(1)', 10, C_GRAY),
    ('8向射线通路(8)  +  目标/左/右开阔度(3)', 10, C_GRAY),
    ('动作前瞻安全余量(1)  +  朝向对齐度(1)', 10, C_GRAY),
])

# Action space
mtext(sl, LX + Inches(0.15), Inches(5.65), LW - Inches(0.3), Inches(1.4), [
    ('▎ 动作空间 (2维) & 环境交互', 14, C_PURPLE, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━', 9, C_LGRAY),
    ('Action[0]=Forward → MoveStepScale映射位移', 10, C_GRAY),
    ('Action[1]=Turn → TurnStepDegrees映射旋转', 10, C_GRAY),
    ('移动: MoveComponent Sweep检测碰撞', 10, C_GRAY),
    ('碰撞后: 追加射线穿透检测+可视化调试', 10, C_GRAY),
    ('bUseBoxSweep=true 盒体扫描(非球体)', 10, C_GRAY),
])

# ═══════════════════════════════════════════════════════════
# RIGHT: 5 Reward Groups
# ═══════════════════════════════════════════════════════════
RX = Inches(4.4)
RW = Inches(8.6)

# Group 1: Path Progress
box(sl, RX, Inches(0.85), Inches(4.25), Inches(1.65), C_PANEL, border_color=C_GREEN)
mtext(sl, RX + Inches(0.15), Inches(0.9), Inches(3.9), Inches(1.55), [
    ('① 路径推进组 (Path Progress)', 14, C_GREEN, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━━━━━', 8, C_LGRAY),
    ('距离缩短奖励:', 11, C_WHITE, True),
    ('  Reward += (PrevDist - CurrDist) × DistanceRewardScale(0.12)', 10, C_GRAY),
    ('  × HeadingGate = Pow(FacingAlignment, 2.0)', 10, C_GRAY),
    ('  → 朝向目标时距离缩短才给高分; 背向目标时门控接近0', 10, C_GRAY),
    ('到达目标奖励:', 11, C_WHITE, True),
    ('  if CurrDist ≤ ReachTargetDistance(50cm): Reward += 10.0', 10, C_GRAY),
])

# Group 2: Direction Alignment
box(sl, RX, Inches(2.65), Inches(4.25), Inches(1.1), C_PANEL, border_color=C_ACCENT)
mtext(sl, RX + Inches(0.15), Inches(2.7), Inches(3.9), Inches(1.0), [
    ('② 方向对齐组 (Direction Alignment)', 14, C_ACCENT, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━', 8, C_LGRAY),
    ('静态对齐: +FacingAlignment × 0.5 × (1 - 0.001×停滞步数)', 10, C_GRAY),
    ('改善奖励: +ΔFacingAlignment × 1.0  (直接奖励有效转向)', 10, C_GRAY),
    ('实际移动方向: +ActualMoveTargetDot × 0.08  (你走向哪里?)', 10, C_GRAY),
    ('阈值奖励: if 对齐从≤0.95→>0.95: +0.2  (突破性对齐bonus)', 10, C_GRAY),
])

# Group 3: Obstacle Awareness
box(sl, RX, Inches(3.9), Inches(4.25), Inches(1.15), C_PANEL, border_color=C_ORANGE)
mtext(sl, RX + Inches(0.15), Inches(3.95), Inches(3.9), Inches(1.05), [
    ('③ 障碍感知组 (Obstacle Awareness)', 14, C_ORANGE, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━', 8, C_LGRAY),
    ('通路开阔度改善: +ΔClearance × 0.2', 10, C_GRAY),
    ('绕障机会改善: +ΔBypassOpportunity × 0.35', 10, C_GRAY),
    ('  BypassOpportunity = max(目标开阔度, 左偏30°开阔度, 右偏30°开阔度)', 9, C_LGRAY),
    ('  使用 EMA平滑 (α=0.35) 降低射线高频噪声', 9, C_LGRAY),
    ('射线突开探索: if 某方向射线增量>0.18 且 与移动方向对齐>0.707', 10, C_GRAY),
    ('  → +RaySuddenOpeningReward(0.06) × 方向对齐度', 10, C_GRAY),
])

# Group 4: Safety & Step Cost
box(sl, RX, Inches(5.2), Inches(4.25), Inches(0.95), C_PANEL, border_color=C_ACCENT2)
mtext(sl, RX + Inches(0.15), Inches(5.25), Inches(3.9), Inches(0.85), [
    ('④ 安全与步成本组 (Safety & Step Cost)', 14, C_ACCENT2, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━', 8, C_LGRAY),
    ('碰撞惩罚: if bHasCollided → -1.5', 10, C_GRAY),
    ('盲目前冲惩罚: if Forward>0.2 且 (背向目标 或 安全余量<0.3) → -0.35×Forward', 10, C_GRAY),
    ('步成本: 每步 -0.02  转向时保留20% (仅 -0.004, 鼓励先转向再前进)', 10, C_GRAY),
])

# Group 5: Anti-Stagnation
box(sl, RX, Inches(6.3), Inches(4.25), Inches(0.9), C_PANEL, border_color=C_PURPLE)
mtext(sl, RX + Inches(0.15), Inches(6.35), Inches(3.9), Inches(0.8), [
    ('⑤ 反停滞组 (Anti-Stagnation)  ← 多层级检测', 14, C_PURPLE, True),
    ('━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━', 8, C_LGRAY),
    ('位置停滞: 位移<1.0cm → 惩罚 0.2+0.05×连续步数  连续25步→截断episode', 10, C_GRAY),
    ('距离停滞: 目标距离改善<2.0cm 连续12步→惩罚 0.15×倍率', 10, C_GRAY),
    ('   两种停滞独立检测，互补覆盖"贴墙不动"和"绕圈不近"两种失败模式', 9, C_LGRAY),
])

# ═══════════════════════════════════════════════════════════
# RIGHT BOTTOM: Summary formula
# ═══════════════════════════════════════════════════════════
box(sl, Inches(8.9), Inches(0.85), Inches(4.1), Inches(0.7), RGBColor(0x1A, 0x22, 0x2A), border_color=C_PURPLE)
mtext(sl, Inches(9.05), Inches(0.9), Inches(3.8), Inches(0.6), [
    ('总奖励 = 路径推进 + 方向对齐 + 障碍感知', 10, C_WHITE, True),
    ('            + 安全与步成本 + 反停滞', 10, C_WHITE, True),
    ('设计意图: 密集塑形(Dense Shaping)使Critic', 9, C_LGRAY),
    ('在每一帧都有明确的学习信号', 9, C_LGRAY),
])

# ═══════════════════════════════════════════════════════════
# FOOTER: Code reference
# ═══════════════════════════════════════════════════════════
txt(sl, Inches(0.5), Inches(7.2), Inches(12), Inches(0.2),
    'XNavigationCubeEnvComponent.cpp:571-700  ComputeReward()  |  :255-285  PerformRaycasts()  |  :453-567  ApplyAction()',
    size=8, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SAVE & MERGE into defense PPT
# ═══════════════════════════════════════════════════════════
out = Path(__file__).resolve().parent.parent / "docs" / "superpowers" / "specs" / "nav_reward_slide.pptx"
prs.save(str(out))
print(f"Saved: {out}")

# Merge into defense PPT at position 10
defense_path = Path(__file__).resolve().parent.parent / "docs" / "superpowers" / "specs" / "defense_ppt.pptx"
prs_def = Presentation(str(defense_path))
prs_nav = Presentation(str(out))

# Replace slide 9 content (index 9, 0-based = slide 10) with nav reward slide
target_slide = prs_def.slides[9]
nav_slide = prs_nav.slides[0]

# Remove existing shapes
for shape in list(target_slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

# Copy background
try:
    target_slide.background.fill.solid()
    target_slide.background.fill.fore_color.rgb = nav_slide.background.fill.fore_color.rgb
except:
    pass

# Copy nav slide shapes
for shape in nav_slide.shapes:
    el = deepcopy(shape._element)
    target_slide.shapes._spTree.append(el)

prs_def.save(str(defense_path))
print(f"Merged into defense PPT at slide 10. Total: {len(prs_def.slides)} slides")
